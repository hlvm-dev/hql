#!/usr/bin/env python3
"""
Generate a 5-slide dark-themed presentation about MCP (Model Context Protocol)
with cyan accents using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import subprocess

# Define colors
DARK_BG = RGBColor(20, 22, 28)  # Very dark blue-gray
CYAN = RGBColor(0, 188, 212)    # Cyan accent
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(230, 230, 230)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_slide_with_title_and_content(prs, title, content_lines, is_code=False):
    """Add a slide with title and bulleted content."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = CYAN
    
    # Add cyan accent line under title
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(0.5), Inches(1.6),
        Inches(2), Inches(0)
    )
    line.line.color.rgb = CYAN
    line.line.width = Pt(4)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(8.6), Inches(4.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    font_size = Pt(20) if is_code else Pt(28)
    
    for i, line_text in enumerate(content_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line_text
        p.font.size = font_size
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0
        if is_code:
            p.font.name = 'Courier New'
    
    return slide

# Slide 1: Intro Slide
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_BG

# Main title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_frame.text = "Model Context Protocol"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(72)
title_para.font.bold = True
title_para.font.color.rgb = CYAN

# Tagline
tagline_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
tagline_frame = tagline_box.text_frame
tagline_frame.text = "Connecting AI Agents to Tools and APIs"
tagline_para = tagline_frame.paragraphs[0]
tagline_para.font.size = Pt(32)
tagline_para.font.italic = True
tagline_para.font.color.rgb = LIGHT_GRAY

# Speaker placeholder
speaker_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
speaker_frame = speaker_box.text_frame
speaker_frame.text = "Presented by: [Speaker Name]"
speaker_para = speaker_frame.paragraphs[0]
speaker_para.font.size = Pt(24)
speaker_para.font.color.rgb = CYAN

# Slide 2: What is MCP
add_slide_with_title_and_content(
    prs,
    "What is MCP?",
    [
        "• Model Context Protocol - a standardized protocol for AI systems",
        "• Enables AI agents to safely invoke tools and access external resources",
        "• Open specification created by Anthropic",
        "• Bridges the gap between LLMs and real-world capabilities"
    ]
)

# Slide 3: How MCP Works (with ASCII diagram)
add_slide_with_title_and_content(
    prs,
    "How MCP Works",
    [
        "┌─────────────┐",
        "│    You      │",
        "└──────┬──────┘",
        "       │ Query/Task",
        "       ▼",
        "┌─────────────────────┐",
        "│   HLVM Agent        │",
        "└──────┬──────────────┘",
        "       │ MCP Request",
        "       ▼",
        "┌─────────────────────┐",
        "│   MCP Server        │",
        "└──────┬──────────────┘",
        "       │ Invoke Tool",
        "       ▼",
        "┌─────────────────────┐",
        "│  Tool / API / DB    │",
        "└─────────────────────┘"
    ],
    is_code=True
)

# Slide 4: Key Benefits
add_slide_with_title_and_content(
    prs,
    "Key Benefits",
    [
        "• Extensibility: Add new tools without modifying core",
        "• Interoperability: Works across different AI platforms",
        "• Safety: Controlled access to resources and capabilities",
        "• Scalability: Compose complex workflows from simple blocks"
    ]
)

# Slide 5: HLVM + MCP Together
add_slide_with_title_and_content(
    prs,
    "HLVM + MCP",
    [
        "• HLVM uses MCP to integrate specialized tools seamlessly",
        "• Agents can invoke filesystem, web, and code tools",
        "• Seamless context passing across tool boundaries",
        "• Unified interface for multi-step AI workflows"
    ]
)

# Slide 6: Get Started
add_slide_with_title_and_content(
    prs,
    "Get Started",
    [
        "1. Install MCP libraries: pip install mcp",
        "2. Define your tools in a Python class",
        "3. Create an MCP Server and register tools",
        "4. Configure your AI agent to connect to MCP Server",
        "5. Start building intelligent workflows!"
    ]
)

# Save presentation
output_path = "/Users/seoksoonjang/dev/hql/mcp_slides.pptx"
prs.save(output_path)
print(f"✓ Presentation saved to {output_path}")

# Open with macOS 'open' command
subprocess.run(["open", output_path], check=True)
print(f"✓ Opened presentation with default viewer")
