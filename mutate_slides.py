#!/usr/bin/env python3
"""
Mutate slides one by one with corrected targets.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import subprocess
import json
import os

STATE_FILE = ".mutation_state"
PPTX_PATH = "/Users/seoksoonjang/dev/hql/mcp_slides.pptx"

CYAN = RGBColor(0, 188, 212)
WHITE = RGBColor(255, 255, 255)

def get_state():
    if os.path.exists(STATE_FILE):
        with open(STATE_FILE, 'r') as f:
            return json.load(f).get('mutation_count', 0)
    return 0

def save_state(count):
    with open(STATE_FILE, 'w') as f:
        json.dump({'mutation_count': count}, f)

def apply_mutation(mutation_num):
    prs = Presentation(PPTX_PATH)
    
    if mutation_num == 1:
        # Already done (Slide 1 title changed)
        print("✓ Mutation 1: Already applied")
    
    elif mutation_num == 2:
        # Slide 3 (index 2): Add "Every tool is a plugin — no code changes needed"
        slide = prs.slides[2]
        found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text
                if "How MCP Works" in text:
                    text_frame = shape.text_frame
                    p = text_frame.add_paragraph()
                    p.text = "• Every tool is a plugin — no code changes needed"
                    p.font.size = Pt(28)
                    p.font.color.rgb = WHITE
                    p.space_before = Pt(12)
                    p.space_after = Pt(12)
                    found = True
                    print("✓ Mutation 2: Added to Slide 3 (How MCP Works)")
                    break
        if not found:
            print("⚠ Mutation 2: Target not found on Slide 3")
    
    elif mutation_num == 3:
        # Slide 4 (index 3): Add "Works with 17+ pre-installed Claude Code plugins"
        slide = prs.slides[3]
        found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text
                if "Key Benefits" in text:
                    text_frame = shape.text_frame
                    p = text_frame.add_paragraph()
                    p.text = "• Works with 17+ pre-installed Claude Code plugins"
                    p.font.size = Pt(28)
                    p.font.color.rgb = WHITE
                    p.space_before = Pt(12)
                    p.space_after = Pt(12)
                    found = True
                    print("✓ Mutation 3: Added to Slide 4 (Key Benefits)")
                    break
        if not found:
            print("⚠ Mutation 3: Target not found on Slide 4")
    
    elif mutation_num == 4:
        # Slide 5 (index 4): Add "One config line unlocks document generation"
        slide = prs.slides[4]
        found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text
                if "HLVM + MCP" in text or "HLVM uses MCP" in text:
                    text_frame = shape.text_frame
                    p = text_frame.add_paragraph()
                    p.text = "• One config line unlocks document generation"
                    p.font.size = Pt(28)
                    p.font.color.rgb = WHITE
                    p.space_before = Pt(12)
                    p.space_after = Pt(12)
                    found = True
                    print("✓ Mutation 4: Added to Slide 5 (HLVM + MCP)")
                    break
        if not found:
            print("⚠ Mutation 4: Target not found on Slide 5")
    
    elif mutation_num == 5:
        # Slide 6 (index 5): Change title to "Start in 60 Seconds"
        slide = prs.slides[5]
        found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text
                if "Get Started" in text:
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = "Start in 60 Seconds"
                    p.font.size = Pt(54)
                    p.font.bold = True
                    p.font.color.rgb = CYAN
                    found = True
                    print("✓ Mutation 5: Updated Slide 6 title")
                    break
        if not found:
            print("⚠ Mutation 5: Target not found on Slide 6")
    
    prs.save(PPTX_PATH)
    subprocess.run(["open", PPTX_PATH], check=True)

# Main
current_state = get_state()
next_mutation = current_state + 1

if next_mutation <= 5:
    print(f"Applying mutation {next_mutation}/5...")
    apply_mutation(next_mutation)
    save_state(next_mutation)
else:
    print("✓ All 5 mutations complete!")
